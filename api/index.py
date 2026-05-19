import os
import json
import hashlib
import urllib.request
import urllib.error
from io import BytesIO

from fastapi import FastAPI, UploadFile, File, Request
from fastapi.responses import JSONResponse
from pypdf import PdfReader
from pptx import Presentation
from supabase import create_client, Client

app = FastAPI()

# Securely load Environment Variables from Vercel
ANTHROPIC_API_KEY = os.environ.get("ANTHROPIC_API_KEY", "")
SUPABASE_URL = os.environ.get("SUPABASE_URL", "")
SUPABASE_KEY = os.environ.get("SUPABASE_KEY", "")

# Initialize Supabase only if keys are present (prevents crashing if not set up yet)
supabase: Client = None
if SUPABASE_URL and SUPABASE_KEY:
    supabase = create_client(SUPABASE_URL, SUPABASE_KEY)

VALID_MODELS = {"claude-haiku-4-5-20251001", "claude-sonnet-4-6", "claude-opus-4-7"}

def get_text_hash(text: str) -> str:
    return hashlib.md5(text.encode('utf-8')).hexdigest()

def extract_text_pptx(data: bytes) -> str:
    prs = Presentation(BytesIO(data))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        parts = [s.text.strip() for s in slide.shapes if hasattr(s, "text") and s.text.strip()]
        if parts:
            slides.append(f"[Slide {i}]\n" + "\n".join(parts))
    return "\n\n".join(slides)

def extract_text_pdf(data: bytes) -> str:
    reader = PdfReader(BytesIO(data))
    pages = []
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text() or ""
        if text.strip():
            pages.append(f"[Page {i}]\n{text.strip()}")
    return "\n\n".join(pages)

def anthropic_request(payload: dict, extra_headers: dict = None):
    body = json.dumps(payload).encode()
    headers = {
        "Content-Type": "application/json",
        "x-api-key": ANTHROPIC_API_KEY,
        "anthropic-version": "2023-06-01",
    }
    if extra_headers:
        headers.update(extra_headers)
    req = urllib.request.Request(
        "https://api.anthropic.com/v1/messages",
        data=body, headers=headers
    )
    with urllib.request.urlopen(req) as resp:
        return json.loads(resp.read())

@app.post("/api/upload")
async def upload_file(file: UploadFile = File(...)):
    try:
        data = await file.read()
        
        # Detect by magic bytes or filename
        if file.filename.lower().endswith('.pdf') or data[:4] == b"%PDF":
            text = extract_text_pdf(data)
        else:
            text = extract_text_pptx(data)

        if not text.strip():
            return JSONResponse(status_code=400, content={"error": "No readable text found in file"})
            
        file_hash = get_text_hash(text)
        
        # Cache text in Supabase
        if supabase:
            db_res = supabase.table("study_materials").select("file_hash").eq("file_hash", file_hash).execute()
            if not db_res.data:
                supabase.table("study_materials").insert({
                    "file_hash": file_hash,
                    "slide_text": text
                }).execute()

        return {"text": text, "file_hash": file_hash}
    except Exception as e:
        return JSONResponse(status_code=500, content={"error": str(e)})

@app.post("/api/ai")
async def handle_ai(request: Request):
    try:
        body = await request.json()
        prompt = body.get("prompt", "")
        model = body.get("model", "claude-haiku-4-5-20251001")
        generation_type = body.get("type") # 'summary', 'glossary', 'flashcards', etc.
        file_hash = body.get("file_hash")
        
        if not ANTHROPIC_API_KEY:
            return JSONResponse(status_code=500, content={"error": "Server missing Anthropic API Key"})
            
        if model not in VALID_MODELS:
            model = "claude-haiku-4-5-20251001"

        # 1. Check Supabase Cache
        if supabase and file_hash and generation_type:
            # Ensure the column exists in your Supabase table for this generation_type
            try:
                db_res = supabase.table("study_materials").select(generation_type).eq("file_hash", file_hash).execute()
                if db_res.data and db_res.data[0].get(generation_type):
                    return {"text": db_res.data[0][generation_type], "cached": True}
            except Exception as db_err:
                print("Cache read error:", db_err) # Fails silently if column doesn't exist yet

        # 2. Generate with AI
        payload = {
            "model": model, "max_tokens": 2048,
            "messages": [{"role": "user", "content": prompt}]
        }
        
        data = anthropic_request(payload)
        generated_text = data["content"][0]["text"]

        # 3. Save back to Supabase
        if supabase and file_hash and generation_type:
            try:
                supabase.table("study_materials").update({
                    generation_type: generated_text
                }).eq("file_hash", file_hash).execute()
            except Exception as db_err:
                print("Cache write error:", db_err)

        return {"text": generated_text, "cached": False}

    except urllib.error.HTTPError as e:
        err = e.read().decode()
        try: msg = json.loads(err).get("error", {}).get("message", err)
        except: msg = err
        return JSONResponse(status_code=e.code, content={"error": msg})
    except Exception as e:
        return JSONResponse(status_code=500, content={"error": str(e)})

@app.post("/api/chat")
async def handle_chat(request: Request):
    try:
        body = await request.json()
        model = body.get("model", "claude-haiku-4-5-20251001")
        slide_text = body.get("slide_text", "")
        history = body.get("history", [])

        if not ANTHROPIC_API_KEY:
            return JSONResponse(status_code=500, content={"error": "Server missing Anthropic API Key"})

        if model not in VALID_MODELS:
            model = "claude-haiku-4-5-20251001"

        system_content = [
            {
                "type": "text",
                "text": f"You are a helpful study assistant. Answer questions based ONLY on the following slide content. Be concise and clear.\n\n{slide_text}",
                "cache_control": {"type": "ephemeral"}
            }
        ]

        messages = [{"role": m["role"], "content": m["content"]} for m in history]

        payload = {
            "model": model,
            "max_tokens": 1024,
            "system": system_content,
            "messages": messages,
        }

        data = anthropic_request(payload, {"anthropic-beta": "prompt-caching-2024-07-31"})
        return {"text": data["content"][0]["text"]}

    except urllib.error.HTTPError as e:
        err = e.read().decode()
        try: msg = json.loads(err).get("error", {}).get("message", err)
        except: msg = err
        return JSONResponse(status_code=e.code, content={"error": msg})
    except Exception as e:
        return JSONResponse(status_code=500, content={"error": str(e)})